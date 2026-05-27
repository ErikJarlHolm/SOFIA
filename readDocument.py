"""
readDocument – Dokumentlesingsmodul for SOFIA
===============================================
Denne modulen leser innholdet fra dokumenter i forskjellige formater:
    - .docx (Word-dokumenter)
    - .pdf (PDF-filer)
    - .txt (ren tekst)
    - .pptx (PowerPoint-presentasjoner)

Eksporterer:
    - READ_DOCUMENT_TOOL_DEFINITION: JSON-definisjon for Foundry function-calling
    - read_document(): Funksjonen som leser og returnerer dokumentinnholdet

Bruk fra agent.py:
    from readDocument import READ_DOCUMENT_TOOL_DEFINITION, read_document
"""

import logging
from pathlib import Path

log = logging.getLogger(__name__)

# ── Mappe der dokumentene som skal klassifiseres ligger ──────────────────────
DOCUMENTS_DIR = Path(r"C:\Users\erikholm\OneDrive - Atea\Documents\Kunder\Atea AI Norge\Dokumentklassifisering")

# ── Verktøydefinisjon for Foundry ────────────────────────────────────────────
# Denne JSON-strukturen registreres i Foundry slik at agenten kan be om
# å lese et dokument. Følger OpenAI function-calling-formatet.

READ_DOCUMENT_TOOL_DEFINITION = {
    "type": "function",
    "name": "read_document",
    "description": (
        "Leser innholdet fra et dokument i dokumentmappen. "
        "Støtter .docx, .pdf, .txt og .pptx-filer. "
        "Returnerer tekstinnholdet fra dokumentet."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Filnavnet på dokumentet som skal leses (f.eks. 'kontrakt.docx').",
            }
        },
        "required": ["filename"],
    },
}


def _read_docx(file_path: Path) -> str:
    """Les tekst fra en .docx-fil ved hjelp av python-docx."""
    from docx import Document
    doc = Document(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _read_pdf(file_path: Path) -> str:
    """Les tekst fra en .pdf-fil ved hjelp av PyPDF2."""
    from PyPDF2 import PdfReader
    reader = PdfReader(str(file_path))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n".join(pages)


def _read_txt(file_path: Path) -> str:
    """Les tekst fra en .txt-fil."""
    return file_path.read_text(encoding="utf-8")


def _read_pptx(file_path: Path) -> str:
    """Les tekst fra en .pptx-fil ved hjelp av python-pptx."""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    slides_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_content.append(text)
        if slide_content:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text)


def read_document(filename: str) -> dict:
    """
    Les innholdet fra et dokument i dokumentmappen.

    Støtter følgende filtyper: .docx, .pdf, .txt, .pptx

    Args:
        filename: Filnavnet (ikke full sti) på dokumentet som skal leses.

    Returnerer:
        {"status": "ok", "filename": "...", "content": "..."} ved suksess
        {"error": "..."} ved feil
    """
    if not filename:
        return {"error": "Ingen filnavn oppgitt. Vennligst oppgi filnavnet på dokumentet."}

    file_path = DOCUMENTS_DIR / filename

    if not file_path.exists():
        # Prøv å finne filen med case-insensitive søk
        matches = [f for f in DOCUMENTS_DIR.iterdir() if f.name.lower() == filename.lower()]
        if matches:
            file_path = matches[0]
        else:
            available = [f.name for f in DOCUMENTS_DIR.iterdir() if f.is_file()]
            return {
                "error": f"Filen '{filename}' ble ikke funnet i dokumentmappen.",
                "available_files": available[:20],  # Begrens til 20 filer i listen
            }

    suffix = file_path.suffix.lower()
    log.info("Leser dokument: %s (type: %s)", file_path.name, suffix)

    try:
        if suffix == ".docx":
            content = _read_docx(file_path)
        elif suffix == ".pdf":
            content = _read_pdf(file_path)
        elif suffix == ".txt":
            content = _read_txt(file_path)
        elif suffix == ".pptx":
            content = _read_pptx(file_path)
        else:
            return {"error": f"Filtypen '{suffix}' støttes ikke. Støttede typer: .docx, .pdf, .txt, .pptx"}

        # Begrens innholdet til maks 15000 tegn for å unngå token-overflow
        if len(content) > 15000:
            content = content[:15000] + "\n\n[... dokumentet er forkortet ...]"

        log.info("Dokument lest: %d tegn", len(content))
        return {"status": "ok", "filename": file_path.name, "content": content}

    except Exception as e:
        log.error("Feil ved lesing av %s: %s", file_path.name, e)
        return {"error": f"Kunne ikke lese filen: {e}"}
